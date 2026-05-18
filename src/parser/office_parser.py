"""
Office Document Parser - supports DOCX, PPTX, XLSX, and legacy/WPS formats
via external conversion (LibreOffice / WPS CLI).
"""

import logging
import os
import subprocess
import sys
import tempfile
import shutil

logger = logging.getLogger(__name__)

def _find_libreoffice():
    """Try to locate LibreOffice executable on the system."""
    candidates = []
    if sys.platform == 'win32':
        candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            r"C:\Program Files\LibreOffice\program\swriter.exe",
            r"C:\Program Files (x86)\LibreOffice\program\swriter.exe",
        ]
    elif sys.platform == 'darwin':
        candidates = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/swriter",
        ]
    else:
        candidates = ["/usr/bin/libreoffice", "/usr/bin/soffice", "/snap/bin/libreoffice"]
    for c in candidates:
        if os.path.isfile(c):
            logger.debug("Found LibreOffice at: %s", c)
            return c
    return None


def _find_wps():
    """Try to locate WPS Office executable on the system."""
    if sys.platform == 'win32':
        for ver in ['12', '11', '10', '9', '8', '7']:
            p = r"C:\Program Files (x86)\WPS Office\{}\wps.exe".format(ver)
            if os.path.isfile(p):
                logger.debug("Found WPS at: %s", p)
                return p
    return None


def _convert_via_libreoffice(src_path, target_ext):
    """Convert a legacy Office file to modern format using LibreOffice CLI."""
    lo_path = _find_libreoffice()
    if not lo_path:
        logger.debug("LibreOffice not found, cannot convert %s", src_path)
        return None
    tmp_dir = tempfile.mkdtemp(prefix="office_conv_")
    try:
        cmd = [lo_path, '--headless', '--convert-to', target_ext.lstrip('.'), '--outdir', tmp_dir, src_path]
        logger.info("Running LibreOffice conversion: %s", ' '.join(cmd))
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode != 0:
            logger.warning("LibreOffice conversion failed for %s: %s", src_path, result.stderr.strip())
            return None
        base = os.path.splitext(os.path.basename(src_path))[0]
        converted = os.path.join(tmp_dir, base + target_ext)
        if os.path.isfile(converted):
            if os.path.getsize(converted) == 0:
                logger.warning("LibreOffice produced empty file for %s (0 bytes)", src_path)
                return None
            logger.info("LibreOffice converted %s -> %s (%d bytes)", src_path, converted, os.path.getsize(converted))
            return converted
        logger.warning("Converted file not found: %s", converted)
        return None
    except subprocess.TimeoutExpired:
        logger.warning("LibreOffice conversion timed out for %s", src_path)
        return None
    except Exception as e:
        logger.exception("LibreOffice conversion error for %s: %s", src_path, e)
        return None
    finally:
        # Clean up temp directory regardless of outcome
        if os.path.isdir(tmp_dir):
            try:
                shutil.rmtree(tmp_dir, ignore_errors=True)
            except Exception as e:
                logger.warning("Failed to clean up LibreOffice temp dir %s: %s", tmp_dir, e)


def _convert_via_wps(src_path, target_ext):
    """Convert a WPS format file using WPS Office CLI (falls back to LibreOffice).

    Note: WPS CLI (/convert) is tried first since it may handle WPS-format files
    (like .wps, .et, .dps) better than LibreOffice. If WPS fails or is unavailable,
    falls back to LibreOffice.
    """
    wps_path = _find_wps()
    if wps_path:
        tmp_dir = tempfile.mkdtemp(prefix="wps_conv_")
        try:
            cmd = [wps_path, '/convert', src_path,
                   '/output', os.path.join(tmp_dir, os.path.basename(src_path).rsplit('.', 1)[0] + target_ext),
                   '/format', target_ext.lstrip('.')]
            logger.info("Running WPS conversion: %s", ' '.join(cmd))
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if result.returncode == 0:
                base = os.path.splitext(os.path.basename(src_path))[0]
                converted = os.path.join(tmp_dir, base + target_ext)
                if os.path.isfile(converted):
                    logger.info("WPS converted %s -> %s", src_path, converted)
                    return converted
            logger.warning("WPS conversion failed for %s", src_path)
        except Exception as e:
            logger.exception("WPS conversion error for %s: %s", src_path, e)
        finally:
            try:
                shutil.rmtree(tmp_dir, ignore_errors=True)
            except Exception as e:
                logger.warning("Failed to clean up WPS temp dir %s: %s", tmp_dir, e)
    else:
        logger.debug("WPS Office not found for %s", src_path)

    # Fall back to LibreOffice if available
    lo_path = _find_libreoffice()
    if lo_path:
        return _convert_via_libreoffice(src_path, target_ext)
    return None


def _auto_convert(src_path, target_ext):
    """Try to convert a legacy/WPS file using any available converter."""
    converted = _convert_via_libreoffice(src_path, target_ext)
    if converted:
        return converted
    return _convert_via_wps(src_path, target_ext)


def parse_office(filepath, filetype, include_tables=False, include_headers_footers=False,
                 include_footnotes=False, annotate_styles=True, max_rows_xlsx=10000,
                 extract_ppt_notes=False):
    """Parse an Office document and return extracted text."""
    actual_filepath = filepath
    actual_filetype = filetype

    # Handle legacy Office formats (.doc, .ppt, .xls)
    legacy_map = {'doc': ('docx', 'MS Word 97-2003'), 'ppt': ('pptx', 'MS PowerPoint 97-2003'),
                  'xls': ('xlsx', 'MS Excel 97-2003')}
    if filetype in legacy_map:
        target_ext, format_name = legacy_map[filetype]
        logger.info("Legacy %s format detected: %s. Converting to %s...",
                    format_name, filepath, target_ext.upper())
        converted = _auto_convert(filepath, '.' + target_ext)
        if converted:
            actual_filepath = converted
            actual_filetype = target_ext
            logger.info("Successfully converted %s -> %s", filepath, actual_filepath)
        else:
            logger.warning("Cannot parse legacy %s file: %s. Install LibreOffice for auto-conversion.",
                           format_name, filepath)
            return {"extract_type": "text",
                    "text": "[Unsupported legacy format: {}]\n[File: {}]\nConsider converting to {} manually.".format(format_name, os.path.basename(filepath), target_ext.upper()),
                    "metadata": {"mime": "application/legacy-office", "note": "Legacy format"}}

    # Handle WPS formats (.wps, .et, .dps)
    wps_map = {'wps': ('docx', 'WPS Writer'), 'et': ('xlsx', 'WPS Spreadsheet'),
               'dps': ('pptx', 'WPS Presentation')}
    if filetype in wps_map:
        target_ext, format_name = wps_map[filetype]
        logger.info("%s format detected: %s. Converting to %s...",
                    format_name, filepath, target_ext.upper())
        converted = _auto_convert(filepath, '.' + target_ext)
        if converted:
            actual_filepath = converted
            actual_filetype = target_ext
            logger.info("Successfully converted %s -> %s", filepath, actual_filepath)
        else:
            logger.warning("Cannot parse %s file: %s. Install LibreOffice or WPS Office.",
                           format_name, filepath)
            return {"extract_type": "text",
                    "text": "[Unsupported format: {}]\n[File: {}]\nInstall LibreOffice or WPS Office.".format(format_name, os.path.basename(filepath)),
                    "metadata": {"mime": "application/wps-office", "note": "WPS format"}}

    # Parse the (possibly converted) modern file
    try:
        if actual_filetype == "docx":
            return _parse_docx(actual_filepath, include_tables, include_headers_footers,
                               include_footnotes, annotate_styles)
        elif actual_filetype == "pptx":
            return _parse_pptx(actual_filepath, extract_ppt_notes)
        elif actual_filetype == "xlsx":
            return _parse_xlsx(actual_filepath, max_rows_xlsx)
        return None
    except Exception as e:
        logger.exception("Failed to parse office file %s (type=%s): %s",
                         actual_filepath, actual_filetype, e)
        return None
    finally:
        if actual_filepath != filepath and os.path.exists(actual_filepath):
            tmp_dir = os.path.dirname(actual_filepath)
            try:
                shutil.rmtree(tmp_dir, ignore_errors=True)
            except Exception as e:
                logger.warning("Failed to clean up temp conversion dir %s: %s", tmp_dir, e)


def _get_style_label(paragraph):
    """Return a human-readable style label for a paragraph."""
    try:
        style = paragraph.style
        if style and style.name:
            name = style.name
            # Both checks merged: 'Heading' prefix OR built-in heading style
            if name.startswith('Heading') or name.startswith('heading') or \
               (hasattr(style, 'builtin') and style.builtin and 'heading' in name.lower()):
                return "[{}]".format(name)
    except Exception:
        pass
    return ""


def _parse_docx(filepath, include_tables=False, include_headers_footers=False,
                include_footnotes=False, annotate_styles=True):
    """Parse a .docx file with configurable extraction."""
    from docx import Document
    doc = Document(filepath)

    text_parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            text_parts.append("")
            continue
        style_label = _get_style_label(para) if annotate_styles else ""
        if style_label and "Heading" in style_label:
            text_parts.append("")
            text_parts.append("=" * 60)
            text_parts.append("{} {}".format(style_label, text))
            text_parts.append("=" * 60)
        else:
            text_parts.append("{}{}".format(style_label + " " if style_label else "", text))

    if include_tables and doc.tables:
        text_parts.append("")
        text_parts.append("--- TABLES ---")
        for ti, table in enumerate(doc.tables, 1):
            rows = []
            for row in table.rows:
                # Deduplicate merged cells: python-docx returns duplicate cell objects for merged cells
                seen_texts = set()
                filtered_cells = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text and cell_text not in seen_texts:
                        seen_texts.add(cell_text)
                        filtered_cells.append(cell_text)
                    elif not cell_text:
                        filtered_cells.append("")
                # If no unique non-empty cells remain, keep at least the first cell
                rt = " | ".join(filtered_cells)
                if rt.strip():
                    rows.append(rt)
            if rows:
                text_parts.append("")
                text_parts.append("[Table %d]" % ti)
                text_parts.extend(rows)

    if include_headers_footers:
        for section in doc.sections:
            if section.header:
                for para in section.header.paragraphs:
                    if para.text.strip():
                        text_parts.append("[Header] {}".format(para.text.strip()))
            if section.footer:
                for para in section.footer.paragraphs:
                    if para.text.strip():
                        text_parts.append("[Footer] {}".format(para.text.strip()))

    if include_footnotes:
        try:
            from docx.oxml.ns import qn
            for rel in doc.part.rels.values():
                if "footnotes" in rel.reltype:
                    fn_part = rel.target_part
                    fn_texts = []
                    for fn in fn_part._element.findall(qn('w:footnote')):
                        txt = ''.join(t.text or '' for t in fn.iter(qn('w:t')))
                        if txt.strip():
                            fn_texts.append(txt.strip())
                    if fn_texts:
                        text_parts.append("")
                        text_parts.append("--- FOOTNOTES ---")
                        text_parts.extend(fn_texts)
                    break
        except Exception as e:
            logger.debug("Could not extract footnotes: %s", e)

    normalized = _normalize_paragraphs(text_parts)
    return {"extract_type": "text", "text": normalized,
            "metadata": {"mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        "format": "docx", "annotated_styles": annotate_styles}}


def _normalize_paragraphs(parts):
    """Collapse excessive blank lines while preserving paragraph structure."""
    result = []
    blank_count = 0
    for part in parts:
        if part == "":
            blank_count += 1
            if blank_count <= 2:
                result.append("")
        else:
            blank_count = 0
            result.append(part)
    while result and result[0] == "":
        result.pop(0)
    while result and result[-1] == "":
        result.pop()
    return "\n".join(result)


def _parse_pptx(filepath, extract_notes=False):
    """Parse a .pptx file."""
    from pptx import Presentation
    prs = Presentation(filepath)
    text_parts = []
    for sidx, slide in enumerate(prs.slides, 1):
        text_parts.append("")
        text_parts.append("-" * 50)
        text_parts.append("Slide %d" % sidx)
        text_parts.append("-" * 50)
        has_content = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        text_parts.append(txt)
                        has_content = True
        if extract_notes and slide.has_notes_slide:
            ns = slide.notes_slide
            if ns and ns.notes_text_frame:
                nt = ns.notes_text_frame.text.strip()
                if nt:
                    text_parts.append("")
                    text_parts.append("[Notes]")
                    text_parts.append(nt)
        if not has_content:
            text_parts.append("[Slide content empty]")
    return {"extract_type": "text", "text": _normalize_paragraphs(text_parts),
            "metadata": {"mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        "format": "pptx"}}


def _parse_xlsx(filepath, max_rows=10000):
    """Parse a .xlsx file with row limit per worksheet."""
    import openpyxl
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text_parts.append("")
        text_parts.append("=" * 60)
        text_parts.append("Sheet: {}".format(sheet_name))
        text_parts.append("=" * 60)
        rc = 0  # Reset per-worksheet counter
        for row in ws.iter_rows(values_only=True):
            rc += 1
            if rc > max_rows:
                text_parts.append("")
                text_parts.append("... [Max %d rows reached; truncating sheet] ..." % max_rows)
                break
            rt = " | ".join(str(c) if c is not None else "" for c in row)
            if rt.strip():
                text_parts.append(rt)
    text = _normalize_paragraphs(text_parts)
    wb.close()
    return {"extract_type": "text", "text": text,
            "metadata": {"mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        "format": "xlsx"}}